import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors - Matches Bangkok Tracker Dark Tech & Emerald Aesthetic
    BG_COLOR = RGBColor(10, 17, 32)       # #0a1120 Deep Dark Blue
    BG_CARD = RGBColor(18, 28, 48)        # #121c30 Glass Panel Dark
    BG_CARD_LIGHT = RGBColor(24, 38, 64)  # #182640 Card highlight
    BG_CODE = RGBColor(6, 11, 22)         # #060b16 Code Box
    EMERALD = RGBColor(34, 197, 94)       # #22c55e Emerald
    EMERALD_DARK = RGBColor(16, 185, 129) # #10b981 Teal Emerald
    CYAN = RGBColor(6, 182, 212)          # #06b6d4 Cyan
    TEXT_WHITE = RGBColor(248, 250, 252)  # #f8fafc White
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8 Slate
    TEXT_CODE = RGBColor(226, 232, 240)   # #e2e8f0 Light Code Text
    BORDER_COLOR = RGBColor(30, 48, 80)   # #1e3050 Border

    def set_slide_background(slide, color=BG_COLOR):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text=None, slide_num=None):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = 'Arial'
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED
            p2.font.name = 'Arial'
            p2.space_before = Pt(4)

        if slide_num is not None:
            foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
            ftf = foot_box.text_frame
            ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
            fp = ftf.paragraphs[0]
            fp.text = 'Bangkok 50 Districts Tracker · Architecture Lecture'
            fp.font.size = Pt(9)
            fp.font.color.rgb = TEXT_MUTED
            fp.font.name = 'Arial'

            num_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.3))
            ntf = num_box.text_frame
            ntf.margin_left = ntf.margin_top = ntf.margin_right = ntf.margin_bottom = 0
            np = ntf.paragraphs[0]
            np.text = str(slide_num)
            np.alignment = PP_ALIGN.RIGHT
            np.font.size = Pt(9)
            np.font.color.rgb = TEXT_MUTED
            np.font.name = 'Arial'

    def add_card(slide, left, top, width, height, bg_color=BG_CARD, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, RGBColor(7, 12, 24))
    
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.8), Inches(1.1), Inches(1.1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(16, 185, 129)
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.text = '{ }'
    cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(26)
    cp.font.bold = True
    cp.font.color.rgb = RGBColor(255, 255, 255)
    
    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(3.5))
    ttf = tbox.text_frame
    ttf.word_wrap = True
    ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
    
    tp1 = ttf.paragraphs[0]
    tp1.text = 'FULL-STACK ARCHITECTURE'
    tp1.font.size = Pt(36)
    tp1.font.bold = True
    tp1.font.color.rgb = TEXT_WHITE
    tp1.font.name = 'Arial'
    
    tp2 = ttf.add_paragraph()
    tp2.text = 'Model · View · Controller — with a Modern Relational Database'
    tp2.font.size = Pt(20)
    tp2.font.color.rgb = EMERALD
    tp2.font.name = 'Arial'
    tp2.space_before = Pt(12)
    
    tp3 = ttf.add_paragraph()
    tp3.text = 'A hands-on case study built around the Bangkok 50 Districts Tracker project (Next.js 14 · React · Drizzle ORM · PostgreSQL)'
    tp3.font.size = Pt(13)
    tp3.font.color.rgb = TEXT_MUTED
    tp3.font.name = 'Arial'
    tp3.space_before = Pt(28)

    # ==========================================
    # SLIDE 2: Learning Objectives
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, 'Learning Objectives', 'By the end of this lecture, you will be able to:', 2)

    objectives = [
        'Explain the Model–View–Controller pattern and why it separates concerns in modern web apps',
        'Trace a request end-to-end through client events, Next.js route handlers, Drizzle ORM, and React views',
        'Read and write TypeScript code that queries and mutates a relational database from a dedicated model layer',
        'Run, extend, and debug a production-ready full-stack GIS tracker in VS Code with Bun & Next.js'
    ]

    for i, obj_text in enumerate(objectives):
        top_pos = Inches(1.8 + i * 1.25)
        add_card(slide2, Inches(0.8), top_pos, Inches(11.7), Inches(1.05), BG_CARD, BORDER_COLOR)
        
        nb = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), top_pos + Inches(0.22), Inches(0.6), Inches(0.6))
        nb.fill.solid()
        nb.fill.fore_color.rgb = EMERALD_DARK
        nb.line.fill.background()
        nbtf = nb.text_frame
        nbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = nbtf.paragraphs[0]
        np.text = str(i + 1)
        np.alignment = PP_ALIGN.CENTER
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = RGBColor(255, 255, 255)
        
        otb = slide2.shapes.add_textbox(Inches(2.0), top_pos + Inches(0.28), Inches(10.2), Inches(0.6))
        otf = otb.text_frame
        otf.word_wrap = True
        otf.margin_left = otf.margin_top = otf.margin_right = otf.margin_bottom = 0
        op = otf.paragraphs[0]
        op.text = obj_text
        op.font.size = Pt(15)
        op.font.bold = True
        op.font.color.rgb = TEXT_WHITE
        op.font.name = 'Arial'

    # ==========================================
    # SLIDE 3: What Is MVC?
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, 'What Is MVC?', 'A design pattern that splits an application into three interconnected parts, each with one job.', 3)

    c1 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), Inches(2.4), Inches(1.8), Inches(1.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(16, 185, 129)
    c1.line.fill.background()
    c1tf = c1.text_frame
    c1tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    c1p = c1tf.paragraphs[0]
    c1p.text = 'MODEL\nData & rules'
    c1p.alignment = PP_ALIGN.CENTER
    c1p.font.size = Pt(14)
    c1p.font.bold = True
    c1p.font.color.rgb = RGBColor(255, 255, 255)

    c2 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.76), Inches(2.4), Inches(1.8), Inches(1.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(6, 182, 212)
    c2.line.fill.background()
    c2tf = c2.text_frame
    c2tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    c2p = c2tf.paragraphs[0]
    c2p.text = 'CONTROLLER\nTraffic cop'
    c2p.alignment = PP_ALIGN.CENTER
    c2p.font.size = Pt(14)
    c2p.font.bold = True
    c2p.font.color.rgb = RGBColor(255, 255, 255)

    c3 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.73), Inches(2.4), Inches(1.8), Inches(1.8))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(59, 130, 246)
    c3.line.fill.background()
    c3tf = c3.text_frame
    c3tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    c3p = c3tf.paragraphs[0]
    c3p.text = 'VIEW\nWhat you see'
    c3p.alignment = PP_ALIGN.CENTER
    c3p.font.size = Pt(14)
    c3p.font.bold = True
    c3p.font.color.rgb = RGBColor(255, 255, 255)

    l1 = slide3.shapes.add_textbox(Inches(3.7), Inches(3.0), Inches(2.0), Inches(0.5))
    l1tf = l1.text_frame
    l1p = l1tf.paragraphs[0]
    l1p.text = 'query / mutate'
    l1p.alignment = PP_ALIGN.CENTER
    l1p.font.size = Pt(11)
    l1p.font.italic = True
    l1p.font.color.rgb = TEXT_MUTED

    l2 = slide3.shapes.add_textbox(Inches(7.6), Inches(3.0), Inches(2.0), Inches(0.5))
    l2tf = l2.text_frame
    l2p = l2tf.paragraphs[0]
    l2p.text = 'request / JSON'
    l2p.alignment = PP_ALIGN.CENTER
    l2p.font.size = Pt(11)
    l2p.font.italic = True
    l2p.font.color.rgb = TEXT_MUTED

    add_card(slide3, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.6), BG_CARD, BORDER_COLOR)
    kitb = slide3.shapes.add_textbox(Inches(1.2), Inches(5.1), Inches(10.9), Inches(1.2))
    kitf = kitb.text_frame
    kitf.word_wrap = True
    kitf.margin_left = kitf.margin_top = kitf.margin_right = kitf.margin_bottom = 0
    
    kp1 = kitf.paragraphs[0]
    kp1.text = 'Key Idea: Strict Separation of Responsibilities'
    kp1.font.size = Pt(15)
    kp1.font.bold = True
    kp1.font.color.rgb = EMERALD
    
    kp2 = kitf.add_paragraph()
    kp2.text = 'The Controller never writes SQL directly or exposes internal schema details, and the View never executes database queries or stores secret credentials. Each layer can be tested, refactored, and scaled independently without breaking the others.'
    kp2.font.size = Pt(13)
    kp2.font.color.rgb = TEXT_MUTED
    kp2.space_before = Pt(6)

    # ==========================================
    # SLIDE 4: Why Separate Concerns?
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, 'Why Separate Concerns?', 'Key software engineering benefits of a cleanly decoupled full-stack architecture.', 4)

    cards_data = [
        ('Maintainability', 'Fix bugs, optimize queries, or redesign UI components without cascading unintended side-effects across other layers.', '↻'),
        ('Reusability & Multi-Storage', 'The same Model/Storage layer can power the public interactive map, admin dashboard, export CLI, and API consumers.', '⚙'),
        ('Team Workflow', 'Frontend engineers work on vector SVG animations while backend engineers optimize relational schemas and SQL queries in parallel.', '⌘'),
        ('Testability & Type Safety', 'Strict TypeScript contracts (FullDistrict, Place, TrackerState) guarantee type safety and effortless isolated unit testing.', '✓')
    ]

    for i, (ctitle, cdesc, cicon) in enumerate(cards_data):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.8 + row * 2.5)
        
        add_card(slide4, left, top, Inches(5.7), Inches(2.2), BG_CARD, BORDER_COLOR)
        
        ib = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.35), top + Inches(0.35), Inches(0.65), Inches(0.65))
        ib.fill.solid()
        ib.fill.fore_color.rgb = BG_CARD_LIGHT
        ib.line.color.rgb = BORDER_COLOR
        ibtf = ib.text_frame
        ibtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ip = ibtf.paragraphs[0]
        ip.text = cicon
        ip.alignment = PP_ALIGN.CENTER
        ip.font.size = Pt(16)
        ip.font.color.rgb = EMERALD
        
        ctb = slide4.shapes.add_textbox(left + Inches(1.2), top + Inches(0.3), Inches(4.2), Inches(1.6))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
        
        cp1 = ctf.paragraphs[0]
        cp1.text = ctitle
        cp1.font.size = Pt(16)
        cp1.font.bold = True
        cp1.font.color.rgb = TEXT_WHITE
        
        cp2 = ctf.add_paragraph()
        cp2.text = cdesc
        cp2.font.size = Pt(12)
        cp2.font.color.rgb = TEXT_MUTED
        cp2.space_before = Pt(6)

    # ==========================================
    # SLIDE 5: The Three Layers, Defined
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, 'The Three Layers, Defined', 'Clear division of roles across data, user interface, and application control flow.', 5)

    layers_data = [
        ('MODEL', 'Owns the data.', 
         'Talks to PostgreSQL via Drizzle ORM. Contains database schemas, table relations, cascading deletes, and data validation. Knows nothing about HTTP headers or JSX rendering.',
         'In this project:\nlib/db/schema.ts\nlib/storage.ts\nlib/districts-data.ts',
         EMERALD_DARK),
        ('VIEW', 'Owns the display.',
         'Renders the interactive SVG vector map, district cards, sliding bottom sheet, and analytics modals. Contains zero database queries and minimum business logic.',
         'In this project:\napp/page.tsx\ncomponents/map/BangkokMap.tsx\ncomponents/DistrictBottomSheet.tsx',
         CYAN),
        ('CONTROLLER', 'Owns the flow.',
         'Receives incoming HTTP requests, verifies HttpOnly JWT admin sessions, validates payloads, delegates to the Model, and invalidates server cache tags.',
         'In this project:\napp/api/districts/route.ts\napp/api/auth/login/route.ts\napp/api/export/route.ts',
         RGBColor(99, 102, 241))
    ]

    for i, (ltitle, ltag, ldesc, lproj, lcolor) in enumerate(layers_data):
        left = Inches(0.8 + i * 4.0)
        top = Inches(1.8)
        
        add_card(slide5, left, top, Inches(3.7), Inches(4.8), BG_CARD, lcolor if i == 0 else BORDER_COLOR)
        
        ltb = slide5.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(3.1), Inches(4.2))
        ltf = ltb.text_frame
        ltf.word_wrap = True
        ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = 0
        
        lp1 = ltf.paragraphs[0]
        lp1.text = ltitle
        lp1.font.size = Pt(18)
        lp1.font.bold = True
        lp1.font.color.rgb = lcolor
        
        lp2 = ltf.add_paragraph()
        lp2.text = ltag
        lp2.font.size = Pt(13)
        lp2.font.italic = True
        lp2.font.color.rgb = TEXT_WHITE
        lp2.space_before = Pt(4)
        
        lp3 = ltf.add_paragraph()
        lp3.text = ldesc
        lp3.font.size = Pt(12)
        lp3.font.color.rgb = TEXT_MUTED
        lp3.space_before = Pt(10)
        
        lp4 = ltf.add_paragraph()
        lp4.text = lproj
        lp4.font.size = Pt(10)
        lp4.font.color.rgb = RGBColor(165, 180, 252) if i == 2 else (RGBColor(147, 197, 253) if i == 1 else RGBColor(167, 243, 208))
        lp4.space_before = Pt(16)

    # ==========================================
    # SLIDE 6: Case Study
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, 'Case Study: The Bangkok 50 Districts Tracker App', 
               'A full-stack GIS web application to explore, track, and log visits across all 50 official districts of Bangkok.', 6)

    pillars = [
        ('Next.js 14', 'CONTROLLER', 'App Router API route handlers, serverless endpoints & cache tagging', RGBColor(59, 130, 246)),
        ('React + Tailwind', 'VIEW', 'Interactive vector SVG map, sliding bottom sheet & responsive UI', RGBColor(6, 182, 212)),
        ('Drizzle ORM', 'MODEL', 'Type-safe PostgreSQL schemas, migrations & cascading relations', RGBColor(16, 185, 129)),
        ('Bun Runtime', 'RUNTIME', 'High-performance JavaScript runtime, bundler & package manager', RGBColor(245, 158, 11))
    ]

    for i, (pname, prole, pdesc, pcolor) in enumerate(pillars):
        left = Inches(0.8 + i * 3.0)
        top = Inches(2.0)
        
        add_card(slide6, left, top, Inches(2.7), Inches(3.2), BG_CARD, BORDER_COLOR)
        
        pi = slide6.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), top + Inches(0.35), Inches(0.8), Inches(0.8))
        pi.fill.solid()
        pi.fill.fore_color.rgb = pcolor
        pi.line.fill.background()
        pitf = pi.text_frame
        pitf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pip = pitf.paragraphs[0]
        pip.text = '</>'
        pip.alignment = PP_ALIGN.CENTER
        pip.font.size = Pt(13)
        pip.font.bold = True
        pip.font.color.rgb = RGBColor(255, 255, 255)
        
        ptb = slide6.shapes.add_textbox(left + Inches(0.2), top + Inches(1.3), Inches(2.3), Inches(1.7))
        ptf = ptb.text_frame
        ptf.word_wrap = True
        ptf.margin_left = ptf.margin_top = ptf.margin_right = ptf.margin_bottom = 0
        
        pp1 = ptf.paragraphs[0]
        pp1.text = pname
        pp1.alignment = PP_ALIGN.CENTER
        pp1.font.size = Pt(15)
        pp1.font.bold = True
        pp1.font.color.rgb = TEXT_WHITE
        
        pp2 = ptf.add_paragraph()
        pp2.text = prole
        pp2.alignment = PP_ALIGN.CENTER
        pp2.font.size = Pt(11)
        pp2.font.bold = True
        pp2.font.color.rgb = pcolor
        pp2.space_before = Pt(2)
        
        pp3 = ptf.add_paragraph()
        pp3.text = pdesc
        pp3.alignment = PP_ALIGN.CENTER
        pp3.font.size = Pt(10)
        pp3.font.color.rgb = TEXT_MUTED
        pp3.space_before = Pt(8)

    add_card(slide6, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.1), BG_CARD, BORDER_COLOR)
    bnb = slide6.shapes.add_textbox(Inches(1.1), Inches(5.65), Inches(11.1), Inches(0.8))
    bntf = bnb.text_frame
    bntf.word_wrap = True
    bntf.margin_left = bntf.margin_top = bntf.margin_right = bntf.margin_bottom = 0
    bnp = bntf.paragraphs[0]
    bnp.text = 'Relational Storage Architecture: District statuses and logged places are stored in normalized relational tables with foreign keys — ensuring ACID integrity, instant lookups, and multi-tier fallback (PostgreSQL / Local JSON).'
    bnp.font.size = Pt(12)
    bnp.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 7: Project Structure
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, 'Project Structure', 'Organizing a modern full-stack Next.js 14 application with clear boundaries.', 7)

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8), BG_CODE, BORDER_COLOR)
    stb = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.4))
    stf = stb.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = 0
    
    tree_text = """bangkok-district-tracker/
|-- app/
|   |-- page.tsx              # VIEW: Map & Dashboard
|   |-- admin/page.tsx        # VIEW: Place Manager
|   `-- api/
|       |-- auth/login/       # CONTROLLER: JWT Auth
|       `-- districts/        # CONTROLLER: CRUD API
|-- components/
|   |-- map/BangkokMap.tsx    # VIEW: SVG Map Engine
|   `-- DistrictBottomSheet.tsx # VIEW: Slide-up Drawer
|-- lib/
|   |-- db/schema.ts          # MODEL: Drizzle Schema
|   |-- storage.ts            # MODEL: Data Access
|   `-- districts-data.ts     # MODEL: Metadata & Stats
`-- data/bangkok-districts.json # DATA: SVG Coordinates"""
    
    sp = stf.paragraphs[0]
    sp.text = tree_text
    sp.font.size = Pt(10.5)
    sp.font.color.rgb = TEXT_CODE
    sp.font.name = 'Consolas'

    add_card(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), BG_CARD, BORDER_COLOR)
    rtb = slide7.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.0), Inches(4.4))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = 0
    
    rp1 = rtf.paragraphs[0]
    rp1.text = 'Why organize it this way?'
    rp1.font.size = Pt(16)
    rp1.font.bold = True
    rp1.font.color.rgb = EMERALD
    
    reasons = [
        'app/api/ isolates HTTP request parsing and authentication from visual components.',
        'components/ encapsulates specialized UI widgets (vector map, bottom drawer, modals).',
        'lib/db/ & lib/storage.ts encapsulates all database queries and schema definitions.',
        'data/ maintains static GIS geo-polygons and starter landmark seed datasets.',
        'lib/types.ts provides shared, single-source-of-truth TypeScript definitions.'
    ]
    
    for r in reasons:
        rp = rtf.add_paragraph()
        rp.text = f'• {r}'
        rp.font.size = Pt(11.5)
        rp.font.color.rgb = TEXT_WHITE
        rp.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: The Model Layer
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, 'The Model Layer', 'lib/db/schema.ts & lib/storage.ts — the only files that define and query data', 8)

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(6.4), Inches(4.8), BG_CODE, BORDER_COLOR)
    mctb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(6.0), Inches(4.4))
    mctf = mctb.text_frame
    mctf.word_wrap = True
    mctf.margin_left = mctf.margin_top = mctf.margin_right = mctf.margin_bottom = 0
    
    model_code = """// lib/db/schema.ts
export const districtStatuses = pgTable('district_statuses', {
  districtId: text('district_id').primaryKey(),
  isVisited: boolean('is_visited').notNull().default(false),
  generalNotes: text('general_notes'),
  updatedAt: timestamp('updated_at').defaultNow().notNull(),
});

export const places = pgTable('places', {
  id: text('id').primaryKey(),
  districtId: text('district_id').notNull()
    .references(() => districtStatuses.districtId, 
      { onDelete: 'cascade' }),
  name: text('name').notNull(),
  category: text('category').notNull().default('Other'),
  visitedDate: text('visited_date'),
  notes: text('notes'),
});"""

    mp = mctf.paragraphs[0]
    mp.text = model_code
    mp.font.size = Pt(10.5)
    mp.font.color.rgb = TEXT_CODE
    mp.font.name = 'Consolas'

    add_card(slide8, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8), BG_CARD, BORDER_COLOR)
    mntb = slide8.shapes.add_textbox(Inches(7.8), Inches(2.0), Inches(4.4), Inches(4.4))
    mntf = mntb.text_frame
    mntf.word_wrap = True
    mntf.margin_left = mntf.margin_top = mntf.margin_right = mntf.margin_bottom = 0
    
    mnp1 = mntf.paragraphs[0]
    mnp1.text = 'Notice'
    mnp1.font.size = Pt(16)
    mnp1.font.bold = True
    mnp1.font.color.rgb = EMERALD
    
    model_notes = [
        ('Type-Safe Schema:', 'Drizzle ORM defines tables with TypeScript inference, preventing SQL syntax and type errors at compile time.'),
        ('Cascading Foreign Keys:', 'The references() constraint with onDelete: \'cascade\' automatically cleans up places if a district status changes.'),
        ('Zero HTTP / UI Logic:', 'No req, res, or HTML anywhere in this layer. It is pure data definition and database queries.')
    ]
    for ntitle, ndesc in model_notes:
        np = mntf.add_paragraph()
        np.text = f'{ntitle} {ndesc}'
        np.font.size = Pt(11.5)
        np.font.color.rgb = TEXT_WHITE
        np.space_before = Pt(12)

    # ==========================================
    # SLIDE 9: The View Layer
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, 'The View Layer', 'components/map/BangkokMap.tsx — turns district data into interactive vector graphics', 9)

    add_card(slide9, Inches(0.8), Inches(1.8), Inches(6.4), Inches(4.8), BG_CODE, BORDER_COLOR)
    vctb = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(6.0), Inches(4.4))
    vctf = vctb.text_frame
    vctf.word_wrap = True
    vctf.margin_left = vctf.margin_top = vctf.margin_right = vctf.margin_bottom = 0
    
    view_code = """// components/map/BangkokMap.tsx
export function BangkokMap({ districts, onSelectDistrict }) {
  return (
    <svg viewBox="0 0 1000 800" className="w-full h-auto">
      {districts.map((district) => {
        const isVisited = district.userData?.isVisited;
        return (
          <path
            key={district.id}
            d={district.svgPath}
            fill={isVisited ? '#22c55e' : '#1e293b'}
            stroke="#060913"
            strokeWidth="1.5"
            onClick={() => onSelectDistrict(district)}
            className="transition-colors cursor-pointer hover:opacity-80"
          />
        );
      })}
    </svg>
  );
}"""

    vp = vctf.paragraphs[0]
    vp.text = view_code
    vp.font.size = Pt(10.5)
    vp.font.color.rgb = TEXT_CODE
    vp.font.name = 'Consolas'

    add_card(slide9, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8), BG_CARD, BORDER_COLOR)
    vntb = slide9.shapes.add_textbox(Inches(7.8), Inches(2.0), Inches(4.4), Inches(4.4))
    vntf = vntb.text_frame
    vntf.word_wrap = True
    vntf.margin_left = vntf.margin_top = vntf.margin_right = vntf.margin_bottom = 0
    
    vnp1 = vntf.paragraphs[0]
    vnp1.text = 'Notice'
    vnp1.font.size = Pt(16)
    vnp1.font.bold = True
    vnp1.font.color.rgb = CYAN
    
    view_notes = [
        ('Pure Declarative UI:', 'The view receives data props from SWR and renders colors dynamically without querying the database.'),
        ('Event Dispatching:', 'onClick triggers onSelectDistrict(district), delegating navigation and state logic back to the parent container.'),
        ('Responsive SVG Matrix:', 'Handles pan, pinch-to-zoom, and touch gestures on the client without triggering backend roundtrips.')
    ]
    for ntitle, ndesc in view_notes:
        np = vntf.add_paragraph()
        np.text = f'{ntitle} {ndesc}'
        np.font.size = Pt(11.5)
        np.font.color.rgb = TEXT_WHITE
        np.space_before = Pt(12)

    # ==========================================
    # SLIDE 10: The Controller Layer
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, 'The Controller Layer', 'app/api/districts/route.ts — request validation, security, and flow control', 10)

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(6.4), Inches(4.8), BG_CODE, BORDER_COLOR)
    cctb = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(6.0), Inches(4.4))
    cctf = cctb.text_frame
    cctf.word_wrap = True
    cctf.margin_left = cctf.margin_top = cctf.margin_right = cctf.margin_bottom = 0
    
    controller_code = """// app/api/districts/route.ts
export async function POST(req: NextRequest) {
  // 1. Verify Admin Authentication via JWT
  const isAdmin = await isAuthenticatedAdmin();
  if (!isAdmin) {
    return NextResponse.json({ error: 'Unauthorized' }, { status: 401 });
  }

  // 2. Parse & Validate Request Body
  const { action, districtId, place } = await req.json();
  if (action === 'add_place') {
    if (!place?.name) {
      return NextResponse.json({ error: 'Place name required' }, { status: 400 });
    }
    // 3. Delegate to Model Layer
    const res = await addPlaceToDistrict(districtId, place);
    
    // 4. Invalidate Server Cache & Return Response
    revalidateTag('districts-state');
    return NextResponse.json({ success: true, state: res.state });
  }
}"""

    cp = cctf.paragraphs[0]
    cp.text = controller_code
    cp.font.size = Pt(10)
    cp.font.color.rgb = TEXT_CODE
    cp.font.name = 'Consolas'

    add_card(slide10, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8), BG_CARD, BORDER_COLOR)
    cntb = slide10.shapes.add_textbox(Inches(7.8), Inches(2.0), Inches(4.4), Inches(4.4))
    cntf = cntb.text_frame
    cntf.word_wrap = True
    cntf.margin_left = cntf.margin_top = cntf.margin_right = cntf.margin_bottom = 0
    
    cnp1 = cntf.paragraphs[0]
    cnp1.text = 'Notice'
    cnp1.font.size = Pt(16)
    cnp1.font.bold = True
    cnp1.font.color.rgb = RGBColor(99, 102, 241)
    
    ctrl_notes = [
        ('Authentication Gatekeeper:', 'Verifies HttpOnly JWT session cookies before allowing state modifications, protecting against unauthorized writes.'),
        ('Validation & Delegation:', 'Validates incoming JSON payloads and delegates storage operations to lib/storage.ts without writing raw SQL.'),
        ('Cache Invalidation:', 'Calls revalidateTag(\'districts-state\') to immediately refresh server-cached responses across all clients.')
    ]
    for ntitle, ndesc in ctrl_notes:
        np = cntf.add_paragraph()
        np.text = f'{ntitle} {ndesc}'
        np.font.size = Pt(11.5)
        np.font.color.rgb = TEXT_WHITE
        np.space_before = Pt(12)

    # ==========================================
    # SLIDE 11: Routes
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, 'Routes: Connecting URLs to Controllers', 
               'Next.js App Router maps HTTP methods and URLs to dedicated Route Handler functions.', 11)

    add_card(slide11, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8), BG_CARD, BORDER_COLOR)
    
    th_box = slide11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.4))
    thf = th_box.text_frame
    thp = thf.paragraphs[0]
    thp.text = f"{'METHOD':<10} {'URL':<22} {'CONTROLLER ACTION':<24} {'PURPOSE'}"
    thp.font.size = Pt(11)
    thp.font.bold = True
    thp.font.color.rgb = EMERALD
    thp.font.name = 'Consolas'

    routes_data = [
        ('GET', '/api/districts', 'GET()', 'Fetch all 50 districts, visited status & stats'),
        ('POST', '/api/districts', 'action: \'toggle_visited\'', 'Toggle district visited status (true / false)'),
        ('POST', '/api/districts', 'action: \'add_place\'', 'Add visited spot (Mall, Temple, Cafe, etc.)'),
        ('POST', '/api/districts', 'action: \'delete_place\'', 'Remove a logged spot from a district'),
        ('POST', '/api/auth/login', 'POST(password)', 'Authenticate admin & issue HttpOnly JWT'),
        ('GET', '/api/export', 'GET()', 'Export full database snapshot as JSON'),
        ('POST', '/api/import', 'POST(backup_json)', 'Restore database state from uploaded JSON'),
        ('POST', '/api/reset', 'POST()', 'Reset database and reseed default landmarks')
    ]

    for i, (rmeth, rurl, ract, rpurp) in enumerate(routes_data):
        row_y = 2.45 + i * 0.48
        rbox = slide11.shapes.add_textbox(Inches(1.0), Inches(row_y), Inches(11.3), Inches(0.4))
        rtf = rbox.text_frame
        rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = 0
        rp = rtf.paragraphs[0]
        rp.text = f"{rmeth:<10} {rurl:<22} {ract:<24} {rpurp}"
        rp.font.size = Pt(11)
        rp.font.color.rgb = TEXT_WHITE if rmeth == 'GET' else RGBColor(147, 197, 253)
        rp.font.name = 'Consolas'

    # ==========================================
    # SLIDE 12: Request Lifecycle
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, 'Request Lifecycle: Adding a Place', 
               'Follow one user request from the browser interface all the way to PostgreSQL and back.', 12)

    lifecycle_steps = [
        ('1', 'User submits \'Add Place\' form in District Bottom Sheet', 'VIEW → POST /api/districts', RGBColor(59, 130, 246)),
        ('2', 'Next.js App Router routes request to the API controller', 'ROUTES → app/api/districts/route.ts', RGBColor(6, 182, 212)),
        ('3', 'Controller verifies JWT auth & validates payload input', 'CONTROLLER → isAuthenticatedAdmin()', RGBColor(168, 85, 247)),
        ('4', 'Model layer executes Drizzle ORM insert into PostgreSQL', 'MODEL → db.insert(places).values(...)', RGBColor(16, 185, 129)),
        ('5', 'Cache revalidated; SWR refetches; district glows emerald', 'CONTROLLER → SWR → VIEW', EMERALD)
    ]

    for i, (snum, sdesc, strans, scolor) in enumerate(lifecycle_steps):
        top_pos = Inches(1.8 + i * 0.95)
        add_card(slide12, Inches(0.8), top_pos, Inches(11.7), Inches(0.8), BG_CARD, BORDER_COLOR)
        
        sb = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), top_pos + Inches(0.15), Inches(0.5), Inches(0.5))
        sb.fill.solid()
        sb.fill.fore_color.rgb = scolor
        sb.line.fill.background()
        sbtf = sb.text_frame
        sbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = sbtf.paragraphs[0]
        sp.text = snum
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.color.rgb = RGBColor(255, 255, 255)
        
        sdtb = slide12.shapes.add_textbox(Inches(1.8), top_pos + Inches(0.2), Inches(6.5), Inches(0.4))
        sdtf = sdtb.text_frame
        sdtf.word_wrap = True
        sdtf.margin_left = sdtf.margin_top = sdtf.margin_right = sdtf.margin_bottom = 0
        sdp = sdtf.paragraphs[0]
        sdp.text = sdesc
        sdp.font.size = Pt(13)
        sdp.font.bold = True
        sdp.font.color.rgb = TEXT_WHITE
        
        sttb = slide12.shapes.add_textbox(Inches(8.5), top_pos + Inches(0.2), Inches(3.8), Inches(0.4))
        sttf = sttb.text_frame
        sttf.word_wrap = True
        sttf.margin_left = sttf.margin_top = sttf.margin_right = sttf.margin_bottom = 0
        stp = sttf.paragraphs[0]
        stp.alignment = PP_ALIGN.RIGHT
        stp.text = strans
        stp.font.size = Pt(11)
        stp.font.color.rgb = scolor
        stp.font.name = 'Consolas'

    # ==========================================
    # SLIDE 13: The Relational Database
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, 'The Relational Database', 
               'PostgreSQL stores district states and places across normalized tables with relational integrity.', 13)

    add_card(slide13, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.3), BG_CARD, BORDER_COLOR)
    dtb = slide13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.4), Inches(3.9))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    dtf.margin_left = dtf.margin_top = dtf.margin_right = dtf.margin_bottom = 0
    
    dp1 = dtf.paragraphs[0]
    dp1.text = 'district_statuses table'
    dp1.font.size = Pt(13)
    dp1.font.bold = True
    dp1.font.color.rgb = EMERALD
    
    schema_desc = """district_id (TEXT, Primary Key)
is_visited  (BOOLEAN, Default: false)
general_notes (TEXT, Optional)
photos      (TEXT JSON, Optional)
updated_at  (TIMESTAMP, Default: NOW())

places table
id          (TEXT, Primary Key)
district_id (TEXT, Foreign Key -> district_statuses)
name        (TEXT, Required)
category    (TEXT, Mall/Temple/Cafe/Park/etc.)
visited_date (TEXT, YYYY-MM-DD)
notes       (TEXT, Optional)"""

    dp2 = dtf.add_paragraph()
    dp2.text = schema_desc
    dp2.font.size = Pt(10)
    dp2.font.color.rgb = TEXT_WHITE
    dp2.font.name = 'Consolas'
    dp2.space_before = Pt(4)

    add_card(slide13, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.3), BG_CODE, BORDER_COLOR)
    ddltb = slide13.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(3.9))
    ddltf = ddltb.text_frame
    ddltf.word_wrap = True
    ddltf.margin_left = ddltf.margin_top = ddltf.margin_right = ddltf.margin_bottom = 0
    
    ddl_code = """CREATE TABLE district_statuses (
  district_id text PRIMARY KEY,
  is_visited boolean NOT NULL DEFAULT false,
  general_notes text,
  updated_at timestamp NOT NULL DEFAULT NOW()
);

CREATE TABLE places (
  id text PRIMARY KEY,
  district_id text NOT NULL 
    REFERENCES district_statuses(district_id) 
    ON DELETE CASCADE,
  name text NOT NULL,
  category text NOT NULL DEFAULT 'Other',
  visited_date text,
  notes text
);"""

    ddlp = ddltf.paragraphs[0]
    ddlp.text = ddl_code
    ddlp.font.size = Pt(10)
    ddlp.font.color.rgb = TEXT_CODE
    ddlp.font.name = 'Consolas'

    ptb = slide13.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5))
    ptf = ptb.text_frame
    ptf.margin_left = ptf.margin_top = ptf.margin_right = ptf.margin_bottom = 0
    pp = ptf.paragraphs[0]
    pp.text = 'Try it yourself: Add a tags table with a many-to-many place_tags junction table. Update lib/storage.ts to support tag filtering.'
    pp.font.size = Pt(11.5)
    pp.font.italic = True
    pp.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 14: Running the App in VS Code
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, 'Running the App in VS Code', 
               'Get the full-stack Bangkok District Tracker development environment up and running in seconds.', 14)

    steps = [
        'Open the bkkdmap / bangkok-district-tracker folder in VS Code',
        'Open a terminal: Terminal → New Terminal',
        'Install project dependencies: bun install (or npm install)',
        'Push database schema to Postgres: bun run db:push',
        'Start development server: bun dev and open http://localhost:3000'
    ]

    for i, s in enumerate(steps):
        top_pos = Inches(1.8 + i * 0.95)
        add_card(slide14, Inches(0.8), top_pos, Inches(7.5), Inches(0.8), BG_CARD, BORDER_COLOR)
        
        sb = slide14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), top_pos + Inches(0.18), Inches(0.45), Inches(0.45))
        sb.fill.solid()
        sb.fill.fore_color.rgb = EMERALD_DARK
        sb.line.fill.background()
        sbtf = sb.text_frame
        sbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = sbtf.paragraphs[0]
        sp.text = str(i + 1)
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(13)
        sp.font.bold = True
        sp.font.color.rgb = RGBColor(255, 255, 255)
        
        sdtb = slide14.shapes.add_textbox(Inches(1.7), top_pos + Inches(0.22), Inches(6.4), Inches(0.4))
        sdtf = sdtb.text_frame
        sdtf.word_wrap = True
        sdtf.margin_left = sdtf.margin_top = sdtf.margin_right = sdtf.margin_bottom = 0
        sdp = sdtf.paragraphs[0]
        sdp.text = s
        sdp.font.size = Pt(12)
        sdp.font.bold = True
        sdp.font.color.rgb = TEXT_WHITE

    add_card(slide14, Inches(8.6), Inches(1.8), Inches(3.9), Inches(4.55), BG_CODE, BORDER_COLOR)
    ttb = slide14.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(3.5), Inches(4.1))
    ttf = ttb.text_frame
    ttf.word_wrap = True
    ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
    
    term_text = """$ bun install
$ bun run db:push
$ bun dev

▲ Next.js 14.2.35
- Local:   http://localhost:3000
- Admin:   http://localhost:3000/admin

✓ Ready in 850ms
✓ Database connected"""

    tp = ttf.paragraphs[0]
    tp.text = term_text
    tp.font.size = Pt(11)
    tp.font.color.rgb = EMERALD
    tp.font.name = 'Consolas'

    ftb = slide14.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
    ftf = ftb.text_frame
    ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
    fp = ftf.paragraphs[0]
    fp.text = 'Zero Setup Required: If no external database is configured, the app automatically falls back to local data/bangkok-tracker-state.json.'
    fp.font.size = Pt(10.5)
    fp.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 15: Practice Exercises
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, 'Practice Exercises', 
               'Extend the Bangkok District Tracker app to deepen your understanding of each architectural layer.', 15)

    exercises = [
        ('GIS Coordinate Pins', 'Add latitude and longitude columns to places and render interactive pinpoint markers inside district SVG boundaries.', '📍', EMERALD_DARK),
        ('Multi-User Passports', 'Extend Model & Controller layers from single-admin to multi-tenant user accounts with personal passport trackers.', '🔐', CYAN),
        ('Faceted Category Search', 'Implement dynamic SQL filtering (GET /api/districts?category=Cafe&zone=Thonburi) in Drizzle ORM.', '🏷️', RGBColor(168, 85, 247)),
        ('Cloud Image Storage', 'Integrate AWS S3 or Cloudflare R2 storage in app/api/upload to support high-res photo albums per place.', '📸', RGBColor(245, 158, 11))
    ]

    for i, (etitle, edesc, eicon, ecolor) in enumerate(exercises):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.8 + row * 2.5)
        
        add_card(slide15, left, top, Inches(5.7), Inches(2.2), BG_CARD, BORDER_COLOR)
        
        ptag = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.3), Inches(2.5), Inches(0.45))
        ptag.fill.solid()
        ptag.fill.fore_color.rgb = ecolor
        ptag.line.fill.background()
        ptagtf = ptag.text_frame
        ptagtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptagtf.paragraphs[0]
        pp.text = f'{eicon} {etitle}'
        pp.alignment = PP_ALIGN.CENTER
        pp.font.size = Pt(11)
        pp.font.bold = True
        pp.font.color.rgb = RGBColor(255, 255, 255)
        
        etb = slide15.shapes.add_textbox(left + Inches(0.3), top + Inches(0.9), Inches(5.1), Inches(1.1))
        etf = etb.text_frame
        etf.word_wrap = True
        etf.margin_left = etf.margin_top = etf.margin_right = etf.margin_bottom = 0
        
        ep = etf.paragraphs[0]
        ep.text = edesc
        ep.font.size = Pt(12)
        ep.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 16: Summary & Discussion
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, RGBColor(7, 12, 24))
    
    chk = slide16.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.2), Inches(0.9), Inches(0.9))
    chk.fill.solid()
    chk.fill.fore_color.rgb = EMERALD_DARK
    chk.line.fill.background()
    chktf = chk.text_frame
    chktf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = chktf.paragraphs[0]
    cp.text = '✓'
    cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = RGBColor(255, 255, 255)
    
    stb = slide16.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.8))
    stf = stb.text_frame
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = 0
    sp = stf.paragraphs[0]
    sp.text = 'Summary'
    sp.font.size = Pt(32)
    sp.font.bold = True
    sp.font.color.rgb = TEXT_WHITE

    summary_bullets = [
        'Full-Stack MVC separates an application into Model (data), View (display), and Controller (flow & security).',
        'The Controller is the single coordinator that manages validation, auth guards, and cache revalidation.',
        'Relational databases enforce data integrity via schemas, foreign keys, and cascading relationships.',
        'The same architectural principles scale from small tracker apps to enterprise geospatial platforms.'
    ]

    btb = slide16.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(2.4))
    btf = btb.text_frame
    btf.word_wrap = True
    btf.margin_left = btf.margin_top = btf.margin_right = btf.margin_bottom = 0
    
    for i, b in enumerate(summary_bullets):
        bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        bp.text = f'• {b}'
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_before = Pt(10)

    qtb = slide16.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0))
    qtf = qtb.text_frame
    qtf.word_wrap = True
    qtf.margin_left = qtf.margin_top = qtf.margin_right = qtf.margin_bottom = 0
    
    qp1 = qtf.paragraphs[0]
    qp1.text = 'Questions?'
    qp1.font.size = Pt(20)
    qp1.font.bold = True
    qp1.font.color.rgb = EMERALD
    
    qp2 = qtf.add_paragraph()
    qp2.text = 'Project repository: bangkok-district-tracker • README.md has full setup, API docs & schema details'
    qp2.font.size = Pt(12)
    qp2.font.color.rgb = TEXT_MUTED
    qp2.space_before = Pt(4)

    output_path = os.path.join('presentation', 'Bangkok-District-Tracker-Architecture-Lecture.pptx')
    prs.save(output_path)
    print(f'Presentation saved successfully to {output_path}')

if __name__ == '__main__':
    create_presentation()
